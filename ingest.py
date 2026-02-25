import os
import argparse
from pathlib import Path
from typing import Optional, List, Dict

# ===== Embeddings + FAISS =====
from langchain_openai import OpenAIEmbeddings
from langchain_community.vectorstores import FAISS
from langchain.text_splitter import RecursiveCharacterTextSplitter

# ===== OpenAI (Whisper) =====
from openai import OpenAI

client = OpenAI()

# ===== Lectores opcionales =====
try:
    from PyPDF2 import PdfReader  # pip install PyPDF2
except Exception:
    PdfReader = None

try:
    import docx  # python-docx
except Exception:
    docx = None

# HTML
try:
    from bs4 import BeautifulSoup  # beautifulsoup4
except Exception:
    BeautifulSoup = None

# PPTX
try:
    from pptx import Presentation  # python-pptx
except Exception:
    Presentation = None

# XLSX
try:
    import openpyxl  # openpyxl
except Exception:
    openpyxl = None


# =========================
# Rutas por defecto
# =========================
BASE_DIR = Path(__file__).resolve().parent
DEFAULT_DOCS_DIR = os.environ.get("DOCS_DIR", str(BASE_DIR / "docs"))
EMBEDDINGS_DIR = os.environ.get("EMBEDDINGS_DIR", str(BASE_DIR / "embeddings"))
TRANSCRIPTS_DIR = os.environ.get(
    "TRANSCRIPTS_DIR", str(BASE_DIR / "transcripts")
)

STT_MODEL = os.environ.get("STT_MODEL", "whisper-1")
STT_LANGUAGE = os.environ.get("STT_LANGUAGE", "es")


# =========================
# Lectores de archivos
# =========================
def read_txt_md(p: Path) -> str:
    """Lee archivos .txt o .md."""
    return p.read_text(encoding="utf-8", errors="ignore")


def read_pdf(p: Path) -> str:
    """Lee un PDF usando PyPDF2."""
    if PdfReader is None:
        raise RuntimeError("PyPDF2 no está instalado. Ejecuta: pip install PyPDF2")
    out: List[str] = []
    reader = PdfReader(str(p))
    for page in reader.pages:
        try:
            text = page.extract_text() or ""
        except Exception:
            text = ""
        if text:
            out.append(text)
    return "\n".join(out)


def read_docx_file(p: Path) -> str:
    """Lee un .docx usando python-docx."""
    if docx is None:
        raise RuntimeError("python-docx no está instalado. Ejecuta: pip install python-docx")
    d = docx.Document(str(p))
    return "\n".join(par.text for par in d.paragraphs)


def read_html_file(p: Path) -> str:
    """Lee un .html y limpia etiquetas."""
    if BeautifulSoup is None:
        raise RuntimeError("beautifulsoup4 no está instalado. Ejecuta: pip install beautifulsoup4")
    html = p.read_text(encoding="utf-8", errors="ignore")
    soup = BeautifulSoup(html, "html.parser")
    return soup.get_text(separator="\n")


def read_pptx_file(p: Path) -> str:
    """Lee un .pptx (texto de las diapositivas)."""
    if Presentation is None:
        raise RuntimeError("python-pptx no está instalado. Ejecuta: pip install python-pptx")
    prs = Presentation(str(p))
    texts: List[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
    return "\n".join(texts)


def read_xlsx_file(p: Path) -> str:
    """Lee un .xlsx concatenando el texto de las celdas."""
    if openpyxl is None:
        raise RuntimeError("openpyxl no está instalado. Ejecuta: pip install openpyxl")
    wb = openpyxl.load_workbook(str(p), data_only=True)
    rows_text: List[str] = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c not in (None, "")]
            if cells:
                rows_text.append(" ".join(cells))
    return "\n".join(rows_text)


def transcript_cache_path(audio_path: Path) -> Path:
    """
    Devuelve la ruta donde se guardará la transcripción .txt
    manteniendo la estructura de carpetas dentro de TRANSCRIPTS_DIR.
    """
    rel = audio_path.relative_to(DEFAULT_DOCS_DIR)
    # Ej: docs/47-XXX/Video.mp4 -> transcripts/47-XXX/Video.mp4.txt
    return Path(TRANSCRIPTS_DIR) / rel.with_suffix(audio_path.suffix + ".txt")


def read_audio_or_video(p: Path) -> str:
    """
    Transcribe un MP3/MP4 usando Whisper. Usa caché en transcripts/.
    """
    # Crear carpeta base de transcripts
    Path(TRANSCRIPTS_DIR).mkdir(parents=True, exist_ok=True)

    cache_path = transcript_cache_path(p)
    cache_path.parent.mkdir(parents=True, exist_ok=True)

    # 1) Si la transcripción ya existe, usarla
    if cache_path.is_file():
        print(f"     [CACHE] Usando transcripción existente: {cache_path}")
        return cache_path.read_text(encoding="utf-8", errors="ignore")

    print(f"     [STT] Transcribiendo audio/vídeo con Whisper: {p.name}")

    try:
        with p.open("rb") as f:
            result = client.audio.transcriptions.create(
                model=STT_MODEL,
                file=f,
                language=STT_LANGUAGE,
                response_format="text",
            )
    except Exception as e:
        print(f"     [ERROR] Falló la transcripción de {p.name}: {e}")
        return ""

    text = result if isinstance(result, str) else str(result)

    # Guardar en caché
    try:
        cache_path.write_text(text, encoding="utf-8")
        print(f"     [STT] Transcripción guardada en: {cache_path}")
    except Exception as e:
        print(f"     [WARN] No se pudo guardar la transcripción en caché: {e}")

    return text


def load_document(p: Path) -> str:
    """Devuelve el texto de un archivo soportado."""
    suffix = p.suffix.lower()

    if suffix in {".txt", ".md"}:
        return read_txt_md(p)
    if suffix == ".pdf":
        return read_pdf(p)
    if suffix == ".docx":
        return read_docx_file(p)
    if suffix in {".html", ".htm"}:
        return read_html_file(p)
    if suffix == ".pptx":
        return read_pptx_file(p)
    if suffix == ".xlsx":
        return read_xlsx_file(p)
    if suffix in {".mp3", ".mp4"}:
        return read_audio_or_video(p)

    # Otros tipos de momento se ignoran
    return ""


# =========================
# Construcción de índices
# =========================
def build_index_for_folder(root: Path, course_id: Optional[str], embeddings: OpenAIEmbeddings) -> None:
    """
    Construye un índice FAISS para todos los documentos de una carpeta.
    Si course_id es None, se crea un índice global en EMBEDDINGS_DIR/faiss_index.
    Si course_id tiene valor, se guarda en EMBEDDINGS_DIR/course_<id>/faiss_index.
    """
    if not root.is_dir():
        print(f"[WARN] La ruta {root} no es una carpeta; se omite.")
        return

    print(f"[INFO] Procesando carpeta: {root}")

    splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=150)

    texts: List[str] = []
    metadatas: List[Dict[str, str]] = []

    # Extensiones soportadas
    exts = {
        ".txt",
        ".md",
        ".pdf",
        ".docx",
        ".html",
        ".htm",
        ".pptx",
        ".xlsx",
        ".mp3",
        ".mp4",
    }

    for path in root.rglob("*"):
        if not path.is_file():
            continue
        if path.suffix.lower() not in exts:
            continue

        rel_path = path.relative_to(root)
        print(f"   - Leyendo archivo: {rel_path}")

        try:
            raw = load_document(path)
        except Exception as e:
            print(f"     [ERROR] No se pudo leer {rel_path}: {e}")
            continue

        if not raw or not raw.strip():
            print(f"     [WARN] Archivo {rel_path} sin texto útil.")
            continue

        chunks = splitter.split_text(raw)
        for ch in chunks:
            # Prepend filename to content so vector search matches queries like "What is in File X?"
            enriched_text = f"Archivo: {path.name}\n{ch}"
            texts.append(enriched_text)
            metadatas.append({
                "source": str(rel_path),
                "course_id": course_id or ""
            })

    if not texts:
        print("[WARN] No se encontraron textos para indexar en esta carpeta.")
        return

    print(f"[INFO] Generando embeddings para {len(texts)} fragmentos...")
    vs = FAISS.from_texts(texts, embeddings, metadatas=metadatas)

    if course_id:
        out_dir = Path(EMBEDDINGS_DIR) / f"course_{course_id}" / "faiss_index"
    else:
        out_dir = Path(EMBEDDINGS_DIR) / "faiss_index"

    out_dir.parent.mkdir(parents=True, exist_ok=True)
    vs.save_local(str(out_dir))
    print(f"[OK] Índice guardado en: {out_dir}\n")


def infer_course_id_from_folder_name(folder_name: str) -> Optional[str]:
    """
    A partir de '47-Economía y Sociedad' devuelve '47'.
    Si no hay dígitos iniciales claros, devuelve None.
    """
    base = folder_name.split("-", 1)[0].strip()
    digits = "".join(ch for ch in base if ch.isdigit())
    return digits or None


# =========================
# CLI
# =========================
def main() -> None:
    parser = argparse.ArgumentParser(
        description="Genera índices FAISS para los documentos del asistente (por curso o global)."
    )
    parser.add_argument(
        "--root",
        type=str,
        default=DEFAULT_DOCS_DIR,
        help=f"Carpeta raíz con documentos (por defecto: {DEFAULT_DOCS_DIR})",
    )
    parser.add_argument(
        "--course-id",
        type=str,
        default=None,
        help="ID de curso específico. Si se omite y --per-course está activo, se infiere por carpeta.",
    )
    parser.add_argument(
        "--per-course",
        action="store_true",
        help="Si se activa, trata cada subcarpeta inmediata de root como un curso distinto.",
    )
    args = parser.parse_args()

    root = Path(args.root).resolve()
    print(f"[INFO] Carpeta raíz de documentos: {root}")
    print(f"[INFO] Directorio de embeddings:   {EMBEDDINGS_DIR}")
    print(f"[INFO] Carpeta de transcripciones: {TRANSCRIPTS_DIR}")
    root.mkdir(parents=True, exist_ok=True)
    Path(EMBEDDINGS_DIR).mkdir(parents=True, exist_ok=True)
    Path(TRANSCRIPTS_DIR).mkdir(parents=True, exist_ok=True)

    embeddings = OpenAIEmbeddings()

    if args.per_course:
        # Modo por curso: cada subcarpeta de root es un curso separado
        for sub in sorted(root.iterdir()):
            if not sub.is_dir():
                continue

            raw_name = sub.name
            course_id = args.course_id

            if course_id is None:
                course_id = infer_course_id_from_folder_name(raw_name)

            if not course_id:
                print(f"[WARN] Carpeta {raw_name} no tiene un id válido antes del '-'; se omite.")
                continue

            print(f"[INFO] Carpeta '{raw_name}' -> course_id='{course_id}'")
            build_index_for_folder(sub, course_id, embeddings)
    else:
        # Modo simple: una sola carpeta -> un índice (course_id dado o índice global)
        build_index_for_folder(root, args.course_id, embeddings)


if __name__ == "__main__":
    main()
