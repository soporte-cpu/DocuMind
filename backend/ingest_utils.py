import os
from pathlib import Path
from typing import Optional, List, Dict
from langchain_openai import OpenAIEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_text_splitters import RecursiveCharacterTextSplitter
from openai import OpenAI
from dotenv import load_dotenv

load_dotenv()

client = OpenAI()

# Lectores opcionales
try:
    from PyPDF2 import PdfReader
except ImportError:
    PdfReader = None

try:
    import docx
except ImportError:
    docx = None

try:
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

BASE_DIR = Path(__file__).resolve().parent.parent
DOCS_DIR = BASE_DIR / "docs"
EMBEDDINGS_DIR = BASE_DIR / "embeddings"
TRANSCRIPTS_DIR = BASE_DIR / "transcripts"

STT_MODEL = "whisper-1"
STT_LANGUAGE = "es"

def load_document(p: Path) -> str:
    suffix = p.suffix.lower()
    if suffix in {".txt", ".md"}:
        return p.read_text(encoding="utf-8", errors="ignore")
    if suffix == ".pdf":
        if PdfReader is None: return ""
        reader = PdfReader(str(p))
        return "\n".join(page.extract_text() or "" for page in reader.pages)
    if suffix == ".docx":
        if docx is None: return ""
        d = docx.Document(str(p))
        return "\n".join(par.text for par in d.paragraphs)
    if suffix in {".html", ".htm"}:
        if BeautifulSoup is None: return ""
        html = p.read_text(encoding="utf-8", errors="ignore")
        return BeautifulSoup(html, "html.parser").get_text(separator="\n")
    if suffix == ".pptx":
        if Presentation is None: return ""
        prs = Presentation(str(p))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
        return "\n".join(texts)
    if suffix == ".xlsx":
        if openpyxl is None: return ""
        wb = openpyxl.load_workbook(str(p), data_only=True)
        rows_text = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c not in (None, "")]
                if cells: rows_text.append(" ".join(cells))
        return "\n".join(rows_text)
    if suffix in {".mp3", ".mp4"}:
        return transcribe_audio(p)
    return ""

def transcribe_audio(p: Path) -> str:
    TRANSCRIPTS_DIR.mkdir(parents=True, exist_ok=True)
    cache_path = TRANSCRIPTS_DIR / (p.name + ".txt")
    if cache_path.is_file():
        return cache_path.read_text(encoding="utf-8", errors="ignore")
    
    try:
        with p.open("rb") as f:
            result = client.audio.transcriptions.create(
                model=STT_MODEL,
                file=f,
                language=STT_LANGUAGE,
                response_format="text",
            )
        text = result if isinstance(result, str) else str(result)
        cache_path.write_text(text, encoding="utf-8")
        return text
    except Exception as e:
        print(f"Error transcribiendo {p.name}: {e}")
        return ""

def update_vector_store():
    """Escanea la carpeta docs y regenera el índice FAISS."""
    if not DOCS_DIR.exists():
        DOCS_DIR.mkdir(parents=True, exist_ok=True)
        return

    splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=250)
    texts: List[str] = []
    metadatas: List[Dict[str, str]] = []

    exts = {".txt", ".md", ".pdf", ".docx", ".html", ".htm", ".pptx", ".xlsx", ".mp3", ".mp4"}

    print(f"[INFO] Iniciando vectorización de {len(list(DOCS_DIR.glob('*')))} archivos en {DOCS_DIR}...")
    for path in DOCS_DIR.rglob("*"):
        if not path.is_file(): continue
        
        # Advertencia especial para .doc (formato viejo)
        if path.suffix.lower() == ".doc":
            print(f"   [!] ADVERTENCIA: {path.name} es un archivo .doc antiguo. Por favor, conviértelo a .docx para que el sistema pueda leerlo.")
            continue

        if path.suffix.lower() in exts:
            print(f"   [+] Procesando: {path.name}")
            try:
                raw = load_document(path)
                if not raw or not raw.strip(): 
                    print(f"   [!] Archivo vacío o no soportado: {path.name}")
                    continue
                area = path.parent.name if path.parent != DOCS_DIR else "General"
                chunks = splitter.split_text(raw)
                for ch in chunks:
                    texts.append(f"Archivo: {path.name}\n{ch}")
                    metadatas.append({
                        "source": path.name,
                        "area": area
                    })
            except Exception as e:
                print(f"   [ERROR] Error procesando {path.name}: {e}")

    if not texts:
        print("[WARN] No se extrajo texto de ningún documento. El índice no se actualizará.")
        return

    print(f"[INFO] Generando embeddings para {len(texts)} fragmentos...")
    embeddings = OpenAIEmbeddings()
    vs = FAISS.from_texts(texts, embeddings, metadatas=metadatas)
    
    # Guardar índice FAISS
    EMBEDDINGS_DIR.mkdir(parents=True, exist_ok=True)
    vs.save_local(str(EMBEDDINGS_DIR / "faiss_index"))
    
    # También guardamos los textos originales para BM25 (Híbrido)
    import pickle
    with open(EMBEDDINGS_DIR / "docs_chunks.pkl", "wb") as f:
        pickle.dump({"texts": texts, "metadatas": metadatas}, f)
        
    print(f"[OK] Índices vectoriales y de texto guardados con éxito en {EMBEDDINGS_DIR}")

def get_hybrid_retriever(area: Optional[str] = None):
    """Carga los índices y devuelve un EnsembleRetriever (FAISS + BM25) con filtro opcional por área."""
    from langchain.retrievers.ensemble import EnsembleRetriever
    from langchain_community.retrievers import BM25Retriever
    from langchain_core.documents import Document
    import pickle

    vs = get_vector_store()
    if not vs: return None
    
    chunks_path = EMBEDDINGS_DIR / "docs_chunks.pkl"
    if not chunks_path.exists():
        search_kwargs = {"k": 8}
        if area: search_kwargs["filter"] = {"area": area}
        return vs.as_retriever(search_kwargs=search_kwargs)

    with open(chunks_path, "rb") as f:
        data = pickle.load(f)
    
    # Filtrar chunks por área si se solicita
    filtered_texts = []
    filtered_metadatas = []
    for t, m in zip(data["texts"], data["metadatas"]):
        if not area or m.get("area") == area:
            filtered_texts.append(t)
            filtered_metadatas.append(m)
            
    if not filtered_texts:
        return vs.as_retriever(search_kwargs={"k": 8})

    # Crear objetos Document para BM25
    docs = [Document(page_content=t, metadata=m) for t, m in zip(filtered_texts, filtered_metadatas)]
    
    # Inicializar retrievers
    faiss_kwargs = {"k": 5}
    if area: faiss_kwargs["filter"] = {"area": area}
    
    faiss_retriever = vs.as_retriever(search_kwargs=faiss_kwargs)
    bm25_retriever = BM25Retriever.from_documents(docs)
    bm25_retriever.k = 5
    
    # Combinar ambos (Búsqueda Híbrida)
    # FAISS es bueno para significado semántico
    # BM25 es excelente para términos técnicos exactos como "pospac"
    ensemble = EnsembleRetriever(
        retrievers=[faiss_retriever, bm25_retriever],
        weights=[0.6, 0.4]
    )
    return ensemble

def get_vector_store():
    index_path = EMBEDDINGS_DIR / "faiss_index"
    if not index_path.exists():
        return None
    return FAISS.load_local(str(index_path), OpenAIEmbeddings(), allow_dangerous_deserialization=True)
